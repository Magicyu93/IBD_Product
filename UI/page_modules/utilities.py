from pptx import Presentation
from pptx.util import Pt, Inches
import streamlit as st
import pandas as pd
import os
import json
import sys
from fpdf import FPDF
from mistletoe import markdown
from langchain.chat_models import AzureChatOpenAI
from langchain.schema import (SystemMessage, HumanMessage)
verbose_txt_sample_path = os.path.join(os.path.dirname(__file__), '..', '..', 'AutoGen', 'backend', 'ui_sample_output')
sys.path.append(verbose_txt_sample_path)

llm = AzureChatOpenAI(
    model="gpt-4o",
    azure_endpoint="https://genai-openai-quantifai.openai.azure.com/",
    openai_api_key="1b31fc4eb58c4879960c46f697d72af6",
    api_version="2024-02-01",
    verbose=False,
    temperature=0.0)

def get_summary(verbose_txt, title, testing=False):

    if testing:
        with open(verbose_txt_sample_path+'/slides_titles_text.json', 'r') as f:
            responses = json.load(f)
        concise_txt = responses.get(title, "None")
    else:
        messages = [SystemMessage(content = f"You are an expert in summarizing documents into a format which is appropriate for a single presentation slide."),
                    HumanMessage(content = f"Please provide a concise summary of the following text in bullet points: {verbose_txt}. Use at most 7 bullet points!")]
        response = llm(messages = messages, temperature = 0)
        concise_txt = response.content
    
    return concise_txt

# Function to save content to a PowerPoint file
def save_content_to_ppt(company_ticker, filename="result/Profiler_Slides.pptx"):
    # Ensure the result directory exists
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    
    # Create a PowerPoint presentation
    path = 'template_barclays.pptx'
    prs = Presentation(path)

    possible_slides = ["Overview", "Financials", "Geographic Mix", "Management",
                       "Recent News", "M&A Profile", "Miscellanea", "Discounted Cash Flow Analysis", "Leveraged Buyout Analysis"]

    for title in possible_slides:
        if st.session_state.get(title, False):

            # Add a slide for each section
            slide_layout_num = 16 if title=='Financials' else 15
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_num])
            title_placeholder = slide.placeholders[0]
            content_placeholder = slide.placeholders[14]

            # Set the title and content
            title_placeholder.text = title
            content_placeholder.text = get_summary(st.session_state[title], title, False)

            content_placeholder.text_frame.margin_left = Inches(0.5)
            content_placeholder.text_frame.margin_top = Inches(0.25)
            content_placeholder.text_frame.margin_right = Inches(0.5)
            content_placeholder.text_frame.margin_bottom = Inches(0.25)  

            for paragraph in content_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.space_before = Pt(10)

            if title == 'Financials':
                if company_ticker=='AAPL':
                    slide.shapes.add_picture("result/apple.png", Inches(6), Inches(1.5), width=Inches(4))
                    df = pd.read_csv("result/apple.csv")
                else:
                    slide.shapes.add_picture("result/google.png", Inches(6), Inches(1.5), width=Inches(4))
                    df = pd.read_csv("result/google.csv")
                df['Market Cap'] = df['Market Cap'].apply(lambda l: l.replace(',', ''))
                df['Market Cap'] = '$' + (df['Market Cap'].astype(float)/1000000000).round(2).astype(str) + 'BB'
                df.drop(columns=['Name'], inplace=True)
                # Create a table
                x, y, cx, cy = Inches(5.75), Inches(4), Inches(5), Inches(1.5)
                shape = slide.shapes.add_table(df.shape[0] + 1, df.shape[1], x, y, cx, cy)
                table = shape.table

                # Add headers
                for j, col in enumerate(df.columns):
                    col_name = table.cell(0, j)
                    col_name.text = col
                    for paragraph in col_name.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(14)

                # Add data
                for i, row in df.iterrows():
                    for j, cell_value in enumerate(row):
                        cell = table.cell(i + 1, j)
                        cell.text = str(cell_value)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(12)          

    # Save the presentation
    prs.save(filename)

# Function to save content to a PowerPoint file
def save_content_to_pdf(company_ticker, filename="result/Profiler_Report.pdf"):

    possible_slides = ["Overview", "Financials", "Geographic Mix", "Management",
                       "Recent News", "M&A Profile", "Miscellanea", "Discounted Cash Flow Analysis", "Leveraged Buyout Analysis"]
    pdf = FPDF()
    pdf.add_page()
    pdf.add_font("dejavu-sans", style="", fname="font/DejaVuSans.ttf")
    pdf.add_font("dejavu-sans", style="b", fname="font/DejaVuSans-Bold.ttf")
    for title in possible_slides:
        if st.session_state.get(title, False):

            html_text = markdown(st.session_state[title])
            html_text = html_text.replace('<h3>', '<h3 style="color: grey;">')
            pdf.set_font(family="dejavu-sans", style="b", size=20)
            pdf.cell(200, 10, txt = title, ln = 1, align = 'C')
            pdf.set_font(family="dejavu-sans", style="", size=12)
            pdf.write_html(html_text)
            if title=='Financials':
                pdf.set_x(40)
                if company_ticker=='AAPL':
                    pdf.cell(40, 10, link=pdf.image("result/apple.png", w=100, h=80), align = 'C')
                    df = pd.read_csv("result/apple.csv")
                else:
                    pdf.cell(40, 10, link=pdf.image("result/google.png", w=100, h=80), align = 'C')
                    df = pd.read_csv("result/google.csv")
                df['Market Cap'] = df['Market Cap'].apply(lambda l: l.replace(',', ''))
                df['Market Cap'] = '$' + (df['Market Cap'].astype(float)/1000000000).round(2).astype(str) + 'BB'
                df.drop(columns=['Name'], inplace=True)

                # Add table header
                pdf.set_font(family="dejavu-sans", style="", size=8)
                pdf.set_x(10)
                for col_name in df.columns:
                    pdf.cell(35, 10, col_name, 1, 0, 'C')
                pdf.ln()

                # Add table rows
                for index, row in df.iterrows():
                    for col_value in row:
                        pdf.cell(35, 10, str(col_value), 1, 0, 'C')
                    pdf.ln()
    # save the pdf with name .pdf
    pdf.output(filename)   